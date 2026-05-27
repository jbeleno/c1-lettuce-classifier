"""Generate ``docs/exports/slides.pptx`` for the C1 sustentation.

Same visual identity as the dashboard (slate-900, five class colors,
JetBrains Mono for numbers) so the deck and the demo feel like the same
product. 16:9, ~12 slides, all in English.

Run with::

    make slides       # builds the .pptx (and PDF if libreoffice is around)
"""
from __future__ import annotations

import json
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ── Tokens (mirror frontend/src/index.css and docs/design-system.md) ─────────
BG          = RGBColor(0x0F, 0x17, 0x2A)
SURFACE_1   = RGBColor(0x1E, 0x29, 0x3B)
SURFACE_INS = RGBColor(0x0B, 0x12, 0x20)
INK         = RGBColor(0xF8, 0xFA, 0xFC)
INK_2       = RGBColor(0xCB, 0xD5, 0xE1)
INK_3       = RGBColor(0x94, 0xA3, 0xB8)
INK_MUTE    = RGBColor(0x64, 0x74, 0x8B)
BORDER      = RGBColor(0x33, 0x41, 0x55)
ACCENT_GO   = RGBColor(0x22, 0xC5, 0x5E)  # Ready (green)
ACCENT_PG   = RGBColor(0xA7, 0x8B, 0xFA)  # germination (violet)
ACCENT_YO   = RGBColor(0xFB, 0xBF, 0x24)  # young (amber)
ACCENT_PD   = RGBColor(0x38, 0xBD, 0xF8)  # pod (sky)
ACCENT_EM   = RGBColor(0x64, 0x74, 0x8B)  # empty_pod (slate)
RING_FOCUS  = RGBColor(0x38, 0xBD, 0xF8)

FONT_SANS = "Inter"
FONT_MONO = "JetBrains Mono"

# 16 : 9 in inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ROOT = Path(__file__).resolve().parent.parent
PNG = ROOT / "docs" / "diagrams" / "png"
MOCK = ROOT / "docs" / "mockups"
MODELS = ROOT / "models_saved"
REPORTS = ROOT / "reports"
OUT_PPTX = ROOT / "docs" / "exports" / "slides.pptx"
OUT_PDF = ROOT / "docs" / "exports" / "slides.pdf"

# Final results (Swin-T leads, ensemble best macro_recall). Loaded from
# the metadata.json files so a re-train automatically updates the deck.
def _load_results() -> dict[str, dict]:
    out = {}
    for child in MODELS.iterdir():
        meta = child / "metadata.json"
        if not meta.exists():
            continue
        with open(meta) as f:
            data = json.load(f)
        cr = data.get("classification_report", {}) or {}
        macro = cr.get("macro avg", {}) or {}
        out[data["name"]] = {
            "test_accuracy":    data.get("test_accuracy"),
            "macro_f1":         macro.get("f1-score"),
            "macro_precision": macro.get("precision"),
            "macro_recall":    macro.get("recall"),
            "best_val_accuracy": data.get("best_val_accuracy"),
        }
    return out

RESULTS = _load_results()


# ─── Tiny helpers ────────────────────────────────────────────────────────────


def _fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _border(shape, rgb: RGBColor, width_pt: float = 0.5) -> None:
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def _add_rect(slide, x, y, w, h, fill: RGBColor):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill(rect, fill)
    return rect


def _text(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    size: int = 20,
    color: RGBColor = INK_2,
    font: str = FONT_SANS,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def _new_slide(prs: Presentation, *, page: int | None = None, total: int | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill(slide.background.fill, BG) if False else None
    # Background as a full-bleed rectangle (slide.background.fill on PPTX gets
    # overridden by themes; a rect is reliable across viewers)
    bg = _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    bg.shadow.inherit = False
    # Top hair-line in the cyan focus color — quiet visual anchor
    _add_rect(slide, 0, 0, SLIDE_W, Emu(2 * 9525), RING_FOCUS)
    # Footer (project + page)
    _text(
        slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
        "c1 · hydroponic lettuce growth-stage classifier · USCO BEINSOF52",
        size=9, color=INK_MUTE, font=FONT_MONO,
    )
    if page is not None and total is not None:
        _text(
            slide, Inches(11.8), Inches(7.05), Inches(1.1), Inches(0.3),
            f"{page:02d} / {total:02d}",
            size=9, color=INK_MUTE, font=FONT_MONO, align=PP_ALIGN.RIGHT,
        )
    return slide


def _header(slide, kicker: str, title: str):
    _text(
        slide, Inches(0.5), Inches(0.45), Inches(12), Inches(0.3),
        kicker.upper(), size=10, color=INK_3, font=FONT_MONO,
    )
    _text(
        slide, Inches(0.5), Inches(0.8), Inches(12), Inches(0.8),
        title, size=32, color=INK, font=FONT_SANS, bold=True,
    )


def _add_image(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def _pct(v: float | None) -> str:
    return f"{v * 100:.2f}%" if v is not None else "—"


# ─── Slides ──────────────────────────────────────────────────────────────────


def slide_title(prs: Presentation, total: int) -> None:
    s = _new_slide(prs, page=1, total=total)
    # Big growth-spine on the left (5 dots, the Ready one filled)
    cx0 = Inches(0.9)
    cy = Inches(3.6)
    d = Inches(0.28)
    gap = Inches(0.55)
    colors = [ACCENT_EM, ACCENT_PG, ACCENT_YO, ACCENT_PD, ACCENT_GO]
    for i in range(5):
        cx = cx0 + (d + gap) * i
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
        if i < 4:
            line = s.shapes.add_connector(1, cx + d, cy + d / 2, cx + d + gap, cy + d / 2)
            line.line.color.rgb = colors[i + 1] if i < 4 else INK_MUTE
            line.line.width = Pt(1.5)
        _fill(dot, colors[i])
        _border(dot, colors[i], 1)
    # Title block
    _text(
        s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4),
        "USCO BEINSOF52 · Artificial Intelligence · Project C1",
        size=12, color=INK_3, font=FONT_MONO,
    )
    _text(
        s, Inches(0.9), Inches(2.0), Inches(11), Inches(1.6),
        "Hydroponic Lettuce\nGrowth-Stage Classifier",
        size=54, color=INK, font=FONT_SANS, bold=True,
    )
    _text(
        s, Inches(0.9), Inches(4.4), Inches(11), Inches(0.5),
        "Five deep-learning backbones · ensemble · FastAPI · React",
        size=18, color=INK_2, font=FONT_SANS,
    )
    _text(
        s, Inches(0.9), Inches(5.5), Inches(11), Inches(0.4),
        "Jesús Beleño · Juan Forero", size=18, color=INK, font=FONT_SANS, bold=True,
    )
    _text(
        s, Inches(0.9), Inches(5.95), Inches(11), Inches(0.4),
        "May 2026", size=14, color=INK_MUTE, font=FONT_MONO,
    )


def slide_problem(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "01 · the decision", "Is this pod ready to harvest?")
    body = (
        "Hydroponic operators face a recurrent perception task: every pod in "
        "a tray is at a different point of the growth timeline, and the right "
        "moment to harvest, re-sow or fertilize is decided pod-by-pod.\n\n"
        "Doing this visually for 1,500+ pods a day is feasible but fatiguing. "
        "The cost of error compounds in both directions — harvest early and "
        "yield drops; harvest late and the plant bolts.\n\n"
        "We trained an image classifier that returns the growth stage of one "
        "pod in under a second and persists every decision so the same farm "
        "can audit and retrain."
    )
    _text(s, Inches(0.5), Inches(2.0), Inches(7.5), Inches(4.5),
          body, size=16, color=INK_2)
    # right: 3 micro-stats
    stats = [
        ("19,721", "labeled pod crops"),
        ("5", "growth stages classified"),
        ("93.86%", "test accuracy (Swin-T)"),
    ]
    for i, (big, small) in enumerate(stats):
        top = Inches(2.3) + Inches(1.4) * i
        _add_rect(s, Inches(9), top, Inches(3.8), Inches(1.2), SURFACE_1)
        _text(s, Inches(9.2), top + Inches(0.18), Inches(3.6), Inches(0.6),
              big, size=28, color=INK, font=FONT_MONO, bold=True)
        _text(s, Inches(9.2), top + Inches(0.78), Inches(3.6), Inches(0.4),
              small, size=11, color=INK_3, font=FONT_MONO)


def slide_dataset(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "02 · dataset", "Roboflow lettuce-pallets · 19,721 crops")
    _add_image(s, REPORTS / "class_distribution.png",
               Inches(0.5), Inches(1.9), w=Inches(7.5))
    # right: facts
    facts = [
        ("source frames",       "1,501"),
        ("crops total",         "19,721"),
        ("train · val · test",  "71.4 % · 14.3 % · 14.3 %"),
        ("split strategy",      "Stratified group-aware (no source-frame leak)"),
        ("class balance",       "germ 33.5 %  →  empty 8.1 %  (4.1×)"),
        ("balancing in loss",   'class_weight = "balanced"'),
    ]
    x = Inches(8.7)
    y = Inches(1.9)
    for label, value in facts:
        _text(s, x, y, Inches(4.2), Inches(0.28),
              label.upper(), size=9, color=INK_3, font=FONT_MONO)
        _text(s, x, y + Inches(0.28), Inches(4.2), Inches(0.4),
              value, size=14, color=INK, font=FONT_MONO)
        y += Inches(0.8)


def slide_backbones(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "03 · backbones", "Five ImageNet-pretrained networks compared")
    rows = [
        ("MobileNetV3-Small", "TensorFlow",  "~2.5 M",   "edge baseline"),
        ("EfficientNet-B0",   "TensorFlow",  "~5.3 M",   "midsize CNN"),
        ("ResNet50",          "TensorFlow",  "~25.6 M",  "canonical CNN"),
        ("ViT-B/16",          "PyTorch",     "~86 M",    "global self-attention"),
        ("Swin-Tiny",         "PyTorch",     "~28 M",    "hierarchical attention"),
    ]
    cols = ["model", "framework", "parameters", "role"]
    col_x = [Inches(0.5), Inches(4.0), Inches(7.0), Inches(9.4)]
    col_w = [Inches(3.5), Inches(2.9), Inches(2.3), Inches(3.6)]
    y = Inches(1.9)
    for i, c in enumerate(cols):
        _text(s, col_x[i], y, col_w[i], Inches(0.3),
              c.upper(), size=10, color=INK_3, font=FONT_MONO)
    y += Inches(0.45)
    _add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.02), BORDER)
    y += Inches(0.15)
    accents = [ACCENT_EM, ACCENT_PG, ACCENT_YO, ACCENT_PD, ACCENT_GO]
    for i, (name, fw, params, role) in enumerate(rows):
        # accent square per row
        _add_rect(s, Inches(0.3), y + Inches(0.16), Inches(0.12), Inches(0.12), accents[i])
        _text(s, col_x[0], y, col_w[0], Inches(0.5),
              name, size=14, color=INK, font=FONT_MONO)
        _text(s, col_x[1], y, col_w[1], Inches(0.5),
              fw, size=14, color=INK_3, font=FONT_MONO)
        _text(s, col_x[2], y, col_w[2], Inches(0.5),
              params, size=14, color=INK_2, font=FONT_MONO)
        _text(s, col_x[3], y, col_w[3], Inches(0.5),
              role, size=14, color=INK_2)
        y += Inches(0.85)
    _text(
        s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
        "Same head topology across the five: backbone → pooling → "
        "dropout (TF) → Dense(5, softmax)  with dtype=float32 on the head "
        "so mixed precision stays numerically stable.",
        size=12, color=INK_3, font=FONT_MONO,
    )


def slide_pipeline(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "04 · pipeline", "End-to-end pipeline")
    # The pipeline diagram is portrait-oriented (tall). Constrain by HEIGHT
    # so it can't push into the title above or the footer below, then center
    # horizontally.
    img_h = Inches(5.2)
    pic = _add_image(
        s, PNG / "05_architecture_pipeline.png",
        Inches(0), Inches(1.75), h=img_h,
    )
    if pic is not None:
        pic.left = Emu(int((SLIDE_W - pic.width) / 2))


def slide_training(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "05 · training", "Two-phase fit for CNNs · end-to-end for transformers")
    bullets_l = [
        ("TF CNNs", "phase 1 — head only", "12 ep · Adam(1e-3)"),
        ("TF CNNs", "phase 2 — top 30 % of backbone", "4 ep · Adam(1e-5)"),
        ("TF CNNs", "BatchNorm always frozen", "preserves running stats"),
        ("Torch", "ViT-B/16 + Swin-T end-to-end", "8 ep · AdamW(3e-5)"),
        ("All",   "two-layer class balancing", "sample-level oversampling (balance.py) + loss-level class_weight"),
        ("All",   "augmentation pipeline", "flip · rotation · brightness · contrast · zoom"),
        ("All",   "mixed precision on CUDA",   "fp16 forward · fp32 softmax/loss"),
        ("All",   "checkpoints on best val_acc",   "ModelCheckpoint + EarlyStopping (patience 3)"),
    ]
    y = Inches(1.95)
    for tag, what, how in bullets_l:
        _add_rect(s, Inches(0.5), y, Inches(0.95), Inches(0.55),
                  SURFACE_1)
        _text(s, Inches(0.55), y + Inches(0.14), Inches(0.9), Inches(0.4),
              tag, size=10, color=INK_3, font=FONT_MONO, align=PP_ALIGN.CENTER)
        _text(s, Inches(1.6), y + Inches(0.05), Inches(7), Inches(0.4),
              what, size=15, color=INK, font=FONT_SANS, bold=True)
        _text(s, Inches(1.6), y + Inches(0.32), Inches(11), Inches(0.4),
              how, size=11, color=INK_3, font=FONT_MONO)
        y += Inches(0.65)


def slide_results(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "06 · results", "Held-out test set comparison")
    ordered = sorted(
        RESULTS.items(),
        key=lambda kv: kv[1]["test_accuracy"] or 0,
        reverse=True,
    )
    cols = ["model", "acc.", "macro f1", "macro recall", "best val"]
    col_x = [Inches(0.5), Inches(5.0), Inches(7.0), Inches(9.0), Inches(11.0)]
    col_w = [Inches(4.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    y = Inches(1.95)
    for i, c in enumerate(cols):
        _text(s, col_x[i], y, col_w[i], Inches(0.3),
              c.upper(), size=10, color=INK_3, font=FONT_MONO,
              align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
    y += Inches(0.4)
    _add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.02), BORDER)
    y += Inches(0.15)
    top_acc = max((r["test_accuracy"] or 0) for _, r in ordered)
    for name, r in ordered:
        is_top = (r["test_accuracy"] or 0) == top_acc
        _text(s, col_x[0], y, col_w[0], Inches(0.4),
              name, size=14, color=INK if is_top else INK_2, font=FONT_MONO,
              bold=is_top)
        for i, key in enumerate(
            ["test_accuracy", "macro_f1", "macro_recall", "best_val_accuracy"], start=1
        ):
            _text(s, col_x[i], y, col_w[i], Inches(0.4),
                  _pct(r.get(key)), size=14,
                  color=INK if is_top else INK_2,
                  font=FONT_MONO, align=PP_ALIGN.RIGHT)
        y += Inches(0.55)

    callout = (
        "Swin-Tiny wins overall accuracy. The ensemble wins macro-recall — "
        "the metric aligned with operational cost (every missed Ready costs "
        "harvest timing; every missed empty_pod costs re-sowing)."
    )
    _add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9),
              SURFACE_1)
    _text(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8),
          callout, size=13, color=INK_2)


def slide_confusion(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "07 · per-class", "Where does Swin-Tiny still confuse?")
    cm_img = MODELS / "swin_t" / "confusion_matrix.png"
    # Square image — constrain by height so it fits between header and footer
    _add_image(s, cm_img, Inches(0.5), Inches(1.9), h=Inches(4.8))
    bullets = [
        ("young → pod",
         "the most expensive confusion. Visually similar mid-stage canopies."),
        ("empty_pod",
         "highest per-class precision. Visually distinct (no plant)."),
        ("Ready",
         "high recall. The operationally critical class is well caught."),
        ("germination",
         "isolated. Smallest class but most distinct early-stage cue."),
    ]
    y = Inches(2.0)
    for tag, body in bullets:
        _text(s, Inches(8.4), y, Inches(4.5), Inches(0.4),
              tag, size=14, color=INK, font=FONT_MONO, bold=True)
        _text(s, Inches(8.4), y + Inches(0.36), Inches(4.5), Inches(0.9),
              body, size=11, color=INK_2)
        y += Inches(1.2)


def slide_system(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "08 · system", "FastAPI · PostgreSQL · React · Vite")
    # 4-endpoint API
    eps = [
        ("POST /predict",  "upload → predict → persist"),
        ("GET /history",   "paginated, label/model filters"),
        ("GET /metrics",   "model-comparison table"),
        ("GET /models",    "what is loadable right now"),
    ]
    y = Inches(2.0)
    _text(s, Inches(0.5), y, Inches(6), Inches(0.4),
          "ENDPOINTS", size=10, color=INK_3, font=FONT_MONO)
    y += Inches(0.4)
    for path, desc in eps:
        _text(s, Inches(0.5), y, Inches(4.0), Inches(0.4),
              path, size=14, color=INK, font=FONT_MONO, bold=True)
        _text(s, Inches(4.6), y, Inches(7.5), Inches(0.4),
              desc, size=14, color=INK_3, font=FONT_SANS)
        y += Inches(0.55)

    # Engineering callouts
    notes = [
        ("Lazy model loading", "TF + Torch share one process · RLock so the ensemble loader doesn't self-deadlock"),
        ("TF pinned to CPU",   "PyTorch keeps the GPU; no Metal/CUDA contention"),
        ("StaticPool SQLite",  "in-memory database in tests · production runs on PostgreSQL"),
        ("Vite proxy",         "/api/* → 127.0.0.1:8000 — no CORS in dev"),
    ]
    y = Inches(2.0)
    _text(s, Inches(7.4), y, Inches(6), Inches(0.4),
          "ENGINEERING", size=10, color=INK_3, font=FONT_MONO)
    y += Inches(0.4)
    for title, body in notes:
        _text(s, Inches(7.4), y, Inches(5.5), Inches(0.3),
              title, size=12, color=INK, font=FONT_SANS, bold=True)
        _text(s, Inches(7.4), y + Inches(0.28), Inches(5.5), Inches(0.5),
              body, size=10, color=INK_3, font=FONT_MONO)
        y += Inches(0.85)


def slide_frontend(prs, page, total):
    """Native pptx mockup of the Predict view. We hand-draw it so it stays
    in the slate-900 design system; the salt-rendered PNG would otherwise
    paste a bright-on-white block onto a dark slide."""
    s = _new_slide(prs, page=page, total=total)
    _header(s, "09 · frontend", "React dashboard · webcam · viewfinder crop")

    # ── Hand-drawn mock on the LEFT ────────────────────────────────────────
    card_x = Inches(0.5)
    card_y = Inches(1.85)
    card_w = Inches(7.0)
    card_h = Inches(4.95)
    card = _add_rect(s, card_x, card_y, card_w, card_h, SURFACE_1)
    _border(card, BORDER, 0.5)
    # accent stripe on top in the predicted-class color
    _add_rect(s, card_x, card_y, card_w, Emu(2 * 9525), ACCENT_PG)
    # Card header (mini nav)
    _text(s, card_x + Inches(0.2), card_y + Inches(0.15), Inches(3), Inches(0.3),
          "● c1/lettuce", size=10, color=INK_2, font=FONT_MONO)
    _text(s, card_x + Inches(3.4), card_y + Inches(0.15), Inches(3.5), Inches(0.3),
          "predict · history · metrics · models", size=10, color=INK_3, font=FONT_MONO)
    # accent underline under "predict"
    _add_rect(s, card_x + Inches(3.4), card_y + Inches(0.43),
              Inches(0.55), Emu(1 * 9525), RING_FOCUS)
    # Viewfinder zone with corner brackets
    vf_x = card_x + Inches(0.4)
    vf_y = card_y + Inches(0.7)
    vf_w = Inches(3.4)
    vf_h = Inches(2.4)
    _add_rect(s, vf_x, vf_y, vf_w, vf_h, SURFACE_INS)
    # 4 corner brackets
    def _br(x, y, dx_dir, dy_dir):
        size = Inches(0.18)
        thick = Emu(2 * 9525)
        _add_rect(s, x, y, size if dx_dir > 0 else thick,
                  thick if dx_dir > 0 else size, RING_FOCUS)
        if dx_dir > 0:
            _add_rect(s, x, y, thick, size, RING_FOCUS)
    _br(vf_x, vf_y, +1, +1)
    _br(vf_x + vf_w - Inches(0.18), vf_y, -1, +1)
    _add_rect(s, vf_x + vf_w - Inches(0.18), vf_y, Inches(0.18), Emu(2 * 9525), RING_FOCUS)
    _add_rect(s, vf_x + vf_w - Emu(2 * 9525), vf_y, Emu(2 * 9525), Inches(0.18), RING_FOCUS)
    _br(vf_x, vf_y + vf_h - Inches(0.18), +1, -1)
    _add_rect(s, vf_x, vf_y + vf_h - Emu(2 * 9525), Inches(0.18), Emu(2 * 9525), RING_FOCUS)
    _br(vf_x + vf_w - Inches(0.18), vf_y + vf_h - Inches(0.18), -1, -1)
    _add_rect(s, vf_x + vf_w - Inches(0.18), vf_y + vf_h - Emu(2 * 9525), Inches(0.18), Emu(2 * 9525), RING_FOCUS)
    _add_rect(s, vf_x + vf_w - Emu(2 * 9525), vf_y + vf_h - Inches(0.18), Emu(2 * 9525), Inches(0.18), RING_FOCUS)
    _text(s, vf_x, vf_y + Inches(1.0), vf_w, Inches(0.4),
          "drag, paste or click", size=10, color=INK_3, font=FONT_MONO,
          align=PP_ALIGN.CENTER)

    # Right column inside the mock: spine + confidence + bars
    rc_x = card_x + Inches(4.0)
    rc_y = card_y + Inches(0.85)
    # mini growth spine
    dot_d = Inches(0.18)
    dot_gap = Inches(0.28)
    dot_colors = [ACCENT_EM, ACCENT_PG, ACCENT_YO, ACCENT_PD, ACCENT_GO]
    for i in range(5):
        cx = rc_x + (dot_d + dot_gap) * i
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, rc_y, dot_d, dot_d)
        if i == 1:   # germination active
            _fill(dot, dot_colors[1])
            _border(dot, dot_colors[1], 1)
        else:
            _fill(dot, BG)
            _border(dot, dot_colors[i] if i < 1 else INK_MUTE, 1)
    _text(s, rc_x, rc_y + Inches(0.32), Inches(2), Inches(0.3),
          "germination", size=11, color=ACCENT_PG, font=FONT_MONO, bold=True)
    _text(s, rc_x, rc_y + Inches(0.62), Inches(2.6), Inches(0.6),
          "98.21%", size=30, color=INK, font=FONT_MONO, bold=True)
    _text(s, rc_x, rc_y + Inches(1.2), Inches(2.6), Inches(0.3),
          "via ensemble_avg · 42 ms", size=9, color=INK_3, font=FONT_MONO)

    # Probability bars (5 rows)
    bars = [
        ("empty_pod",   0.00007, ACCENT_EM),
        ("germination", 0.98208, ACCENT_PG),
        ("young",       0.01767, ACCENT_YO),
        ("pod",         0.00015, ACCENT_PD),
        ("Ready",       0.00004, ACCENT_GO),
    ]
    by = rc_y + Inches(1.7)
    max_bar = Inches(1.6)
    for label, p, color in bars:
        _text(s, rc_x, by, Inches(1.0), Inches(0.22),
              label, size=8, color=INK_3 if p < 0.5 else color, font=FONT_MONO)
        track = _add_rect(s, rc_x + Inches(1.05), by + Inches(0.05),
                          max_bar, Emu(2 * 9525), SURFACE_INS)
        if p > 0.001:
            w = Emu(int(max_bar * max(p, 0.005)))
            _add_rect(s, rc_x + Inches(1.05), by + Inches(0.05), w,
                      Emu(2 * 9525), color if p > 0.5 else INK_MUTE)
        _text(s, rc_x + Inches(1.05) + max_bar + Inches(0.08), by, Inches(0.6), Inches(0.22),
              f"{p:.3f}", size=8, color=INK_2 if p > 0.5 else INK_3, font=FONT_MONO,
              align=PP_ALIGN.RIGHT)
        by += Inches(0.28)

    # ── Bullets on the RIGHT ───────────────────────────────────────────────
    bullets = [
        "Vanilla getUserMedia for the laptop webcam — no camera library",
        "Drag-and-drop, clipboard paste or click on the viewfinder",
        "Crop step between picking and predicting — one pod per request",
        "The GrowthSpine signature runs across all four routes",
        "Single accent per card = the predicted class color, never two",
    ]
    y = Inches(2.0)
    for b in bullets:
        _add_rect(s, Inches(8.0), y + Inches(0.08), Emu(2 * 9525), Inches(0.4),
                  RING_FOCUS)
        _text(s, Inches(8.2), y, Inches(4.9), Inches(0.6),
              b, size=12, color=INK_2)
        y += Inches(0.9)


def slide_bug(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "10 · what almost shipped wrong", "A loss of 1,191 on the very first MobileNet run")
    _text(s, Inches(0.5), Inches(2.0), Inches(7.5), Inches(0.4),
          "The smell test", size=14, color=INK_3, font=FONT_MONO)
    body = (
        "MobileNetV3 was passed include_preprocessing=False with no manual\n"
        "normalisation, while modern Keras preprocess_input for MobileNetV3\n"
        "is a no-op. The network was fed pixels in [0, 255] when its ImageNet\n"
        "weights expected [-1, 1].\n\n"
        "Loss exploded into the thousands; test accuracy capped at 0.35 —\n"
        "barely above the 0.20 random baseline. The fix was one flag."
    )
    _text(s, Inches(0.5), Inches(2.5), Inches(7.5), Inches(4),
          body, size=13, color=INK_2, font=FONT_MONO)
    # before / after card
    _add_rect(s, Inches(8.5), Inches(2.0), Inches(4.4), Inches(4.5), SURFACE_1)
    _text(s, Inches(8.7), Inches(2.15), Inches(4), Inches(0.4),
          "BEFORE", size=10, color=INK_3, font=FONT_MONO)
    _text(s, Inches(8.7), Inches(2.45), Inches(4), Inches(0.7),
          "test_acc 0.3465\nloss epoch 1: 1190.88",
          size=18, color=ACCENT_EM, font=FONT_MONO, bold=True)
    _text(s, Inches(8.7), Inches(4.2), Inches(4), Inches(0.4),
          "AFTER", size=10, color=INK_3, font=FONT_MONO)
    _text(s, Inches(8.7), Inches(4.5), Inches(4), Inches(0.7),
          "test_acc 0.8832\nloss epoch 1: 0.59",
          size=18, color=ACCENT_GO, font=FONT_MONO, bold=True)
    _text(s, Inches(8.7), Inches(5.6), Inches(4), Inches(0.6),
          "One line of code · 4 min training\n+0.53 accuracy",
          size=11, color=INK_3, font=FONT_MONO)


def slide_future(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    _header(s, "11 · what's next", "Engineering and product directions")
    pillars = [
        (
            "Detection layer",
            "Add YOLOv8 in front of the classifier. The Roboflow source is "
            "annotated with boxes already — fine-tune on the same data, "
            "then run our 5-model classifier on each detected box. Closes "
            "the multi-pod gap end-to-end.",
        ),
        (
            "Weighted ensemble + stacking",
            "Uniform soft-voting underweights the transformers. Weighted "
            "soft voting (val-accuracy proportional) or a logistic-regression "
            "meta-learner trained on out-of-fold softmax vectors.",
        ),
        (
            "K-fold CV in the report",
            "make cv-swin already exists and writes reports/cv_swin_t/. "
            "Attach a mean ± std to the headline 0.9386 to silence the "
            '"lucky split" question.',
        ),
        (
            "Cloud deployment",
            "Dockerize the FastAPI service + Postgres, push to Render free "
            "tier, get a public https URL the operator opens on a tablet.",
        ),
    ]
    y = Inches(1.95)
    for title, body in pillars:
        _add_rect(s, Inches(0.5), y, Inches(0.08), Inches(0.95), RING_FOCUS)
        _text(s, Inches(0.75), y, Inches(11.8), Inches(0.4),
              title, size=15, color=INK, font=FONT_SANS, bold=True)
        _text(s, Inches(0.75), y + Inches(0.35), Inches(11.8), Inches(0.7),
              body, size=12, color=INK_2)
        y += Inches(1.18)


def slide_thanks(prs, page, total):
    s = _new_slide(prs, page=page, total=total)
    # Center the gratitude
    _text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.5),
          "QUESTIONS · DISCUSSION", size=12, color=INK_3,
          font=FONT_MONO, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.4),
          "Thanks.", size=72, color=INK, font=FONT_SANS, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
          "github.com/jbeleno/c1-lettuce-classifier",
          size=15, color=RING_FOCUS, font=FONT_MONO, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
          "Jesús Beleño · Juan Forero", size=13, color=INK_2,
          font=FONT_SANS, align=PP_ALIGN.CENTER)


# ─── Build + export ──────────────────────────────────────────────────────────


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        slide_title,
        slide_problem,
        slide_dataset,
        slide_backbones,
        slide_pipeline,
        slide_training,
        slide_results,
        slide_confusion,
        slide_system,
        slide_frontend,
        slide_bug,
        slide_future,
        slide_thanks,
    ]
    total = len(slides)
    slides[0](prs, total)
    for page, fn in enumerate(slides[1:], start=2):
        fn(prs, page, total)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"wrote {OUT_PPTX.relative_to(ROOT)}  ({OUT_PPTX.stat().st_size // 1024} KB)")
    return OUT_PPTX


def export_pdf() -> Path | None:
    """Try LibreOffice headless to render a PDF next to the pptx. Skipped
    gracefully if `soffice` isn't on PATH."""
    cmd = ["soffice", "--headless", "--convert-to", "pdf",
           "--outdir", str(OUT_PPTX.parent), str(OUT_PPTX)]
    try:
        subprocess.run(cmd, check=True, capture_output=True, timeout=60)
    except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
        print("libreoffice unavailable — skip PDF export")
        return None
    print(f"wrote {OUT_PDF.relative_to(ROOT)}")
    return OUT_PDF


if __name__ == "__main__":
    build()
    export_pdf()
